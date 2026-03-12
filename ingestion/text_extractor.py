import io
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from typing import Union


def extract_text(file_bytes: bytes, extension: str) -> str:
    """
    Extracts full text from supported document formats.

    Args:
        file_bytes: Raw bytes of the document
        extension: File type — 'pdf', 'docx', 'pptx', 'txt'

    Returns:
        full_document_text as a single string
    """
    if extension == "pdf":
        return _extract_pdf(file_bytes)
    elif extension == "docx":
        return _extract_docx(file_bytes)
    elif extension == "pptx":
        return _extract_pptx(file_bytes)
    elif extension == "txt":
        return _extract_txt(file_bytes)
    else:
        raise ValueError(f"[TextExtractor] Unsupported file format: {extension}")


def _extract_pdf(file_bytes: bytes) -> str:
    """Extracts text from PDF preserving page structure."""
    text_parts = []
    pdf_stream = io.BytesIO(file_bytes)

    # ✅ Use context manager — auto-closes safely after use
    with fitz.open(stream=pdf_stream, filetype="pdf") as doc:
        total_pages = len(doc)
        for page_num, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            if page_text:
                text_parts.append(f"[PAGE {page_num}]\n{page_text}")

    full_text = "\n\n".join(text_parts)
    print(f"[TextExtractor] PDF extracted — {total_pages} pages")
    return full_text

# def _extract_pdf(file_bytes: bytes) -> str:
#     """Extracts text from PDF preserving page structure."""
#     text_parts = []
#     pdf_stream = io.BytesIO(file_bytes)
#     doc = fitz.open(stream=pdf_stream, filetype="pdf")

#     for page_num, page in enumerate(doc, start=1):
#         page_text = page.get_text("text").strip()
#         if page_text:
#             text_parts.append(f"[PAGE {page_num}]\n{page_text}")

#     doc.close()
#     full_text = "\n\n".join(text_parts)
#     print(f"[TextExtractor] PDF extracted — {len(doc)} pages")
#     return full_text


def _extract_docx(file_bytes: bytes) -> str:
    """Extracts text from DOCX paragraph by paragraph."""
    docx_stream = io.BytesIO(file_bytes)
    doc = DocxDocument(docx_stream)
    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
    full_text = "\n\n".join(paragraphs)
    print(f"[TextExtractor] DOCX extracted — {len(paragraphs)} paragraphs")
    return full_text


def _extract_pptx(file_bytes: bytes) -> str:
    """Extracts text from PPTX slide by slide."""
    pptx_stream = io.BytesIO(file_bytes)
    prs = Presentation(pptx_stream)
    slide_texts = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_content.append(shape.text.strip())
        if slide_content:
            slide_texts.append(f"[SLIDE {slide_num}]\n" + "\n".join(slide_content))

    full_text = "\n\n".join(slide_texts)
    print(f"[TextExtractor] PPTX extracted — {len(prs.slides)} slides")
    return full_text


def _extract_txt(file_bytes: bytes) -> str:
    """Extracts text from plain TXT file."""
    full_text = file_bytes.decode("utf-8", errors="ignore").strip()
    print(f"[TextExtractor] TXT extracted — {len(full_text)} characters")
    return full_text
